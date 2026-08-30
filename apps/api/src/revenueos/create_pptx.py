from __future__ import annotations

import io
import re
import zipfile
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import PurePosixPath
from typing import Final, cast
from xml.etree import ElementTree

from pptx import Presentation

PPTX_MIME_TYPE: Final = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PPTX_PROCESSING_SCHEMA_VERSION: Final = 1
PPTX_RENDERER_VERSION: Final = "deterministic_pptx_v1"


class PptxProcessingError(Exception):
    code = "pptx_processing_failed"


class UnsupportedPptxError(PptxProcessingError):
    code = "unsupported_pptx"


class MalformedPptxError(PptxProcessingError):
    code = "malformed_pptx"


class UnsafePptxError(PptxProcessingError):
    code = "unsafe_pptx"


class PptxLimitError(PptxProcessingError):
    code = "pptx_limit_exceeded"


@dataclass(frozen=True)
class PptxLimits:
    max_bytes: int
    max_slides: int
    max_entries: int
    max_expanded_bytes: int
    max_media_assets: int
    max_media_bytes: int
    max_xml_bytes: int
    max_extracted_characters: int


@dataclass(frozen=True)
class ParsedTextBlock:
    shape_id: int
    shape_name: str
    text: str
    placeholder_type: str | None
    editable: bool
    mapped_role: str | None


@dataclass(frozen=True)
class ParsedTemplateSlide:
    slide_number: int
    title: str
    category: str
    reuse_state: str
    modification_policy: str
    customer_safe: bool
    required: bool
    exact_text_required: bool
    hidden: bool
    text_blocks: tuple[ParsedTextBlock, ...]


@dataclass(frozen=True)
class ParsedTemplate:
    width_emu: int
    height_emu: int
    slides: tuple[ParsedTemplateSlide, ...]
    warning_codes: tuple[str, ...]


@dataclass(frozen=True)
class RenderSlide:
    slide_number: int
    replacements: dict[int, str]


_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_FORBIDDEN_RELATIONSHIP_MARKERS = (
    "/oleobject",
    "/control",
    "/activex",
    "/vbaproject",
    "/externallink",
    "/embeddedfont",
)
_FORBIDDEN_PATH_MARKERS = (
    "vbaproject",
    "/activex/",
    "/embeddings/",
    "/oleobjects/",
    "/packages/",
    "/fonts/",
    "/externalLinks/",
    "/customxml/",
)
_ALLOWED_MEDIA_SIGNATURES: tuple[tuple[bytes, str], ...] = (
    (b"\x89PNG\r\n\x1a\n", "png"),
    (b"\xff\xd8\xff", "jpeg"),
    (b"GIF87a", "gif"),
    (b"GIF89a", "gif"),
)
_INTERNAL_ONLY_TERMS = (
    "win probability",
    "deal probability",
    "forecast category",
    "manager coaching",
    "private note",
    "internal risk",
    "contactability",
    "suppression",
    "champion hypothesis",
    "competitive trap",
)


class BoundedPptxProcessor:
    """Validates untrusted PPTX ZIP/XML and extracts a bounded structural manifest."""

    def __init__(self, limits: PptxLimits) -> None:
        self.limits = limits

    def parse(self, content: bytes) -> ParsedTemplate:
        entries = self._preflight(content)
        try:
            presentation = Presentation(io.BytesIO(content))
        except Exception as exc:
            raise MalformedPptxError from exc
        if len(presentation.slides) == 0:
            raise MalformedPptxError
        if len(presentation.slides) > self.limits.max_slides:
            raise PptxLimitError
        hidden_slides = self._hidden_slide_numbers(entries)
        warnings: set[str] = set()
        if any(path.startswith("ppt/notesSlides/") for path in entries):
            warnings.add("speaker_notes_excluded")
        if any(path.startswith("ppt/comments/") or path == "ppt/commentAuthors.xml" for path in entries):
            warnings.add("comments_excluded")
        if hidden_slides:
            warnings.add("hidden_slides_excluded")

        character_count = 0
        parsed_slides: list[ParsedTemplateSlide] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text_blocks: list[ParsedTextBlock] = []
            all_text: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = _normalise_text(shape.text)
                if not text:
                    continue
                character_count += len(text)
                if character_count > self.limits.max_extracted_characters:
                    raise PptxLimitError
                placeholder_type = _placeholder_type(shape)
                mapped_role = _default_placeholder_role(slide_number, placeholder_type)
                text_blocks.append(
                    ParsedTextBlock(
                        shape_id=shape.shape_id,
                        shape_name=shape.name[:200],
                        text=text,
                        placeholder_type=placeholder_type,
                        editable=mapped_role is not None,
                        mapped_role=mapped_role,
                    )
                )
                all_text.append(text)
            title = _slide_title(slide_number, text_blocks)
            category = _classify_slide(slide_number, title, "\n".join(all_text))
            hidden = slide_number in hidden_slides
            internal_only = _contains_internal_only_text("\n".join(all_text))
            customer_safe = not hidden and not internal_only and category != "pricing_placeholder"
            reuse_state = "excluded" if not customer_safe else "pending"
            has_editable_placeholder = any(block.mapped_role is not None for block in text_blocks)
            modification_policy = "text_placeholders_only" if has_editable_placeholder else "reuse_as_is"
            parsed_slides.append(
                ParsedTemplateSlide(
                    slide_number=slide_number,
                    title=title,
                    category=category,
                    reuse_state=reuse_state,
                    modification_policy=modification_policy,
                    customer_safe=customer_safe,
                    required=False,
                    exact_text_required=False,
                    hidden=hidden,
                    text_blocks=tuple(text_blocks),
                )
            )
        width = presentation.slide_width
        height = presentation.slide_height
        if width is None or height is None:
            raise MalformedPptxError
        return ParsedTemplate(
            width_emu=int(width),
            height_emu=int(height),
            slides=tuple(parsed_slides),
            warning_codes=tuple(sorted(warnings)),
        )

    def render(
        self,
        source: bytes,
        slides: tuple[RenderSlide, ...],
        *,
        title: str,
        organisation_name: str,
    ) -> bytes:
        self._preflight(source)
        if not slides or len(slides) > 30:
            raise PptxLimitError
        slide_numbers = [item.slide_number for item in slides]
        if len(slide_numbers) != len(set(slide_numbers)):
            raise UnsafePptxError
        try:
            presentation = Presentation(io.BytesIO(source))
            source_slides = list(presentation.slides)
            if any(number < 1 or number > len(source_slides) for number in slide_numbers):
                raise UnsafePptxError
            selected = [source_slides[number - 1] for number in slide_numbers]
            for render_slide, slide in zip(slides, selected, strict=True):
                replacements = render_slide.replacements
                for shape in slide.shapes:
                    replacement = replacements.get(shape.shape_id)
                    if replacement is None:
                        continue
                    if not shape.has_text_frame:
                        raise UnsafePptxError
                    shape.text = _normalise_text(replacement)
            slide_id_list = presentation.slides._sldIdLst  # noqa: SLF001 - python-pptx exposes no public reorder API
            source_slide_ids = list(slide_id_list)
            selected_slide_ids = [source_slide_ids[number - 1] for number in slide_numbers]
            for slide_id_element in source_slide_ids:
                slide_id_list.remove(slide_id_element)
                if slide_id_element not in selected_slide_ids:
                    presentation.part.drop_rel(slide_id_element.rId)
            for slide_id_element in selected_slide_ids:
                slide_id_list.append(slide_id_element)
            del organisation_name
            presentation.core_properties.title = title[:255]
            presentation.core_properties.subject = "Customer presentation"
            presentation.core_properties.author = "RevenueOS"
            presentation.core_properties.last_modified_by = "RevenueOS"
            presentation.core_properties.comments = ""
            presentation.core_properties.keywords = ""
            presentation.core_properties.category = ""
            presentation.core_properties.content_status = ""
            presentation.core_properties.identifier = ""
            presentation.core_properties.language = ""
            presentation.core_properties.version = ""
            presentation.core_properties.revision = 1
            generated_at = datetime.now(UTC).replace(tzinfo=None)
            presentation.core_properties.created = generated_at
            presentation.core_properties.modified = generated_at
            output = io.BytesIO()
            presentation.save(output)
        except PptxProcessingError:
            raise
        except Exception as exc:
            raise PptxProcessingError from exc
        sanitised = _strip_internal_parts(output.getvalue(), selected_slide_count=len(slides))
        self._preflight(sanitised)
        return sanitised

    def _preflight(self, content: bytes) -> dict[str, bytes]:
        if not content or len(content) > self.limits.max_bytes:
            raise PptxLimitError
        if not content.startswith(b"PK") or not zipfile.is_zipfile(io.BytesIO(content)):
            raise MalformedPptxError
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as archive:
                infos = archive.infolist()
                if len(infos) > self.limits.max_entries:
                    raise PptxLimitError
                expanded = 0
                media_count = 0
                names: set[str] = set()
                entries: dict[str, bytes] = {}
                for info in infos:
                    name = _safe_zip_name(info.filename)
                    folded = name.casefold()
                    if folded in names:
                        raise UnsafePptxError
                    names.add(folded)
                    if info.flag_bits & 0x1:
                        raise UnsafePptxError
                    if info.compress_type not in {zipfile.ZIP_STORED, zipfile.ZIP_DEFLATED}:
                        raise UnsupportedPptxError
                    expanded += info.file_size
                    if expanded > self.limits.max_expanded_bytes:
                        raise PptxLimitError
                    if info.file_size > 1_000_000 and info.compress_size > 0:
                        if info.file_size / info.compress_size > 200:
                            raise PptxLimitError
                    if any(marker.casefold() in f"/{folded}" for marker in _FORBIDDEN_PATH_MARKERS):
                        raise UnsafePptxError
                    if info.file_size > max(self.limits.max_xml_bytes, self.limits.max_media_bytes):
                        raise PptxLimitError
                    data = archive.read(info)
                    if len(data) != info.file_size:
                        raise MalformedPptxError
                    if name.startswith("ppt/media/"):
                        media_count += 1
                        if media_count > self.limits.max_media_assets or len(data) > self.limits.max_media_bytes:
                            raise PptxLimitError
                        _validate_media(data)
                    if name.endswith((".xml", ".rels")):
                        if len(data) > self.limits.max_xml_bytes:
                            raise PptxLimitError
                        _parse_safe_xml(data)
                    entries[name] = data
        except PptxProcessingError:
            raise
        except (OSError, RuntimeError, zipfile.BadZipFile, zipfile.LargeZipFile) as exc:
            raise MalformedPptxError from exc
        required = {"[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"}
        if not required.issubset(entries):
            raise MalformedPptxError
        slide_count = len([name for name in entries if re.fullmatch(r"ppt/slides/slide[1-9][0-9]*\.xml", name)])
        if slide_count == 0:
            raise MalformedPptxError
        if slide_count > self.limits.max_slides:
            raise PptxLimitError
        _validate_content_types(entries["[Content_Types].xml"])
        for name, data in entries.items():
            if not name.endswith(".rels"):
                continue
            root = _parse_safe_xml(data)
            for relationship in root.findall(f"{{{_REL_NS}}}Relationship"):
                if relationship.attrib.get("TargetMode", "").casefold() == "external":
                    raise UnsafePptxError
                relationship_type = relationship.attrib.get("Type", "").casefold()
                if any(marker in relationship_type for marker in _FORBIDDEN_RELATIONSHIP_MARKERS):
                    raise UnsafePptxError
        return entries

    @staticmethod
    def _hidden_slide_numbers(entries: dict[str, bytes]) -> set[int]:
        root = _parse_safe_xml(entries["ppt/presentation.xml"])
        slide_ids = root.findall(f".//{{{_P_NS}}}sldId")
        return {index for index, item in enumerate(slide_ids, start=1) if item.attrib.get("show") == "0"}


def _safe_zip_name(value: str) -> str:
    if not value or "\x00" in value or "\\" in value or value.startswith(("/", ".")):
        raise UnsafePptxError
    path = PurePosixPath(value)
    if any(part in {"", ".", ".."} for part in path.parts):
        raise UnsafePptxError
    normalised = str(path)
    if normalised != value.rstrip("/") and not value.endswith("/"):
        raise UnsafePptxError
    return normalised


def _parse_safe_xml(content: bytes) -> ElementTree.Element:
    lowered = content[:4096].lower()
    if b"<!doctype" in lowered or b"<!entity" in lowered:
        raise UnsafePptxError
    try:
        return ElementTree.fromstring(content)
    except ElementTree.ParseError as exc:
        raise MalformedPptxError from exc


def _serialise_open_xml(root: ElementTree.Element) -> bytes:
    """Preserve the default root namespace required by strict OOXML readers."""
    namespace = root.tag.partition("}")[0].removeprefix("{") if root.tag.startswith("{") else ""
    if namespace:
        ElementTree.register_namespace("", namespace)
    return cast(bytes, ElementTree.tostring(root, encoding="utf-8", xml_declaration=True))


def _validate_content_types(content: bytes) -> None:
    root = _parse_safe_xml(content)
    types = [item.attrib.get("ContentType", "").casefold() for item in root]
    if not any("presentationml.presentation.main+xml" in value for value in types):
        raise UnsupportedPptxError
    if any("macroenabled" in value or "vbaproject" in value for value in types):
        raise UnsafePptxError


def _validate_media(content: bytes) -> None:
    if not content or not any(content.startswith(signature) for signature, _ in _ALLOWED_MEDIA_SIGNATURES):
        raise UnsupportedPptxError


def _normalise_text(value: str) -> str:
    return "\n".join(line.rstrip() for line in value.replace("\r\n", "\n").replace("\r", "\n").splitlines()).strip()


def _placeholder_type(shape: object) -> str | None:
    is_placeholder = bool(getattr(shape, "is_placeholder", False))
    if not is_placeholder:
        return None
    placeholder_format = getattr(shape, "placeholder_format", None)
    placeholder_value = getattr(placeholder_format, "type", None)
    return str(placeholder_value) if placeholder_value is not None else None


def _default_placeholder_role(slide_number: int, placeholder_type: str | None) -> str | None:
    if placeholder_type is None:
        return None
    folded = placeholder_type.casefold()
    if "title" in folded and "subtitle" not in folded:
        return "presentation_title"
    if "subtitle" in folded:
        return "account_name"
    if any(key in folded for key in ("body", "object", "text")):
        return "customer_context" if slide_number > 1 else "audience"
    return None


def _slide_title(slide_number: int, blocks: list[ParsedTextBlock]) -> str:
    for block in blocks:
        if (
            block.placeholder_type
            and "title" in block.placeholder_type.casefold()
            and "subtitle" not in block.placeholder_type.casefold()
        ):
            return block.text.splitlines()[0][:240]
    if blocks:
        return blocks[0].text.splitlines()[0][:240]
    return f"Slide {slide_number}"


def _classify_slide(slide_number: int, title: str, all_text: str) -> str:
    folded = f"{title}\n{all_text}".casefold()
    rules = (
        ("agenda", ("agenda", "today's discussion", "today’s discussion")),
        ("case_study", ("case study", "customer story")),
        ("proof_point", ("proof point", "results", "outcomes")),
        ("architecture", ("architecture", "technical design", "integration")),
        ("next_steps", ("next step", "what happens next")),
        ("process", ("implementation", "our process", "delivery approach")),
        ("company_overview", ("about us", "who we are", "company overview")),
        ("problem", ("challenge", "what we understand", "priorities", "your needs")),
        ("solution", ("solution", "proposed approach", "our approach")),
        ("capability", ("capability", "capabilities")),
        ("product", ("product", "platform")),
        ("pricing_placeholder", ("pricing", "price", "commercials")),
        ("appendix", ("appendix", "disclaimer", "legal notice", "terms")),
    )
    if slide_number == 1:
        return "title"
    for category, keywords in rules:
        if any(keyword in folded for keyword in keywords):
            return category
    return "unknown"


def _contains_internal_only_text(value: str) -> bool:
    folded = value.casefold()
    return any(term in folded for term in _INTERNAL_ONLY_TERMS)


def _strip_internal_parts(content: bytes, *, selected_slide_count: int) -> bytes:
    removed_prefixes = (
        "ppt/notesSlides/",
        "ppt/notesMasters/",
        "ppt/comments/",
        "ppt/commentAuthors.xml",
        "customXml/",
        "docProps/custom.xml",
        "docProps/thumbnail.jpeg",
    )
    source = io.BytesIO(content)
    target = io.BytesIO()
    with zipfile.ZipFile(source) as archive, zipfile.ZipFile(target, "w", compression=zipfile.ZIP_DEFLATED) as output:
        for info in archive.infolist():
            name = info.filename
            if any(name.startswith(prefix) for prefix in removed_prefixes):
                continue
            data = archive.read(info)
            if name.endswith(".rels"):
                root = _parse_safe_xml(data)
                for relationship in list(root):
                    relationship_type = relationship.attrib.get("Type", "").casefold()
                    target_name = relationship.attrib.get("Target", "").casefold()
                    if any(
                        marker in relationship_type or marker.strip("/") in target_name
                        for marker in (
                            "/notesslide",
                            "/notesmaster",
                            "/comments",
                            "/commentauthors",
                            "/thumbnail",
                        )
                    ):
                        root.remove(relationship)
                data = _serialise_open_xml(root)
            elif name == "[Content_Types].xml":
                root = _parse_safe_xml(data)
                for item in list(root):
                    part_name = item.attrib.get("PartName", "").lstrip("/")
                    if any(part_name.startswith(prefix) for prefix in removed_prefixes):
                        root.remove(item)
                data = _serialise_open_xml(root)
            elif name == "docProps/app.xml":
                root = _parse_safe_xml(data)
                for item in list(root):
                    local_name = item.tag.rsplit("}", 1)[-1]
                    if local_name in {
                        "Company",
                        "HeadingPairs",
                        "HyperlinkBase",
                        "Manager",
                        "Template",
                        "TitlesOfParts",
                    }:
                        root.remove(item)
                    elif local_name == "Slides":
                        item.text = str(selected_slide_count)
                data = _serialise_open_xml(root)
            output.writestr(name, data)
    return target.getvalue()
