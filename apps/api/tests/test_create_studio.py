from __future__ import annotations

import asyncio
import base64
import concurrent.futures
import hashlib
import io
import uuid
import zipfile
from dataclasses import replace
from datetime import UTC, datetime, timedelta

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from sqlalchemy import update
from sqlalchemy.ext.asyncio import async_sessionmaker, create_async_engine

from revenueos.auth import AuthenticatedUser, get_current_user
from revenueos.beta_maintenance import run_retention
from revenueos.create_pptx import (
    BoundedPptxProcessor,
    GeneratedOutputValidationError,
    MalformedPptxError,
    OutputSlideExpectation,
    PptxLimitError,
    PptxLimits,
    RenderSlide,
    UnsafePptxError,
    UnsupportedPptxError,
)
from revenueos.create_worker import CreateWorkerService, _validate_claim_manifest, create_processor
from revenueos.models import CreateDownloadGrant, CreatePresentationVersion, OrganisationMembership
from revenueos.visual_storage import create_visual_storage

from .conftest import PRIMARY_ORGANISATION_ID, PRIMARY_USER_ID, TEST_DB_URL, TEST_VISUAL_STORAGE
from .test_business_api import create_company, create_contact, create_opportunity
from .test_business_cases import _create_approved_model, _inputs
from .test_meeting_api import cast_auth_dependency, secondary_user


def _pptx() -> bytes:
    presentation = Presentation()
    title = presentation.slides.add_slide(presentation.slide_layouts[0])
    title.shapes.title.text = "Summit Access Systems"
    title.placeholders[1].text = "Approved company presentation"
    synthetic_logo = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9YEksEIAAAAASUVORK5CYII="
    )
    title.shapes.add_picture(io.BytesIO(synthetic_logo), Inches(11.4), Inches(0.35), Inches(0.55))
    for heading, body in (
        ("Agenda", "Customer context\nApproved capabilities\nProposed approach\nNext steps"),
        (
            "Product platform",
            "Centralised access administration\nPolicy-based controls\nCustomer-owned implementation checkpoints",
        ),
        (
            "Approved customer story",
            "A synthetic customer consolidated access workflows using the approved Summit delivery process.",
        ),
        ("Proposed solution", "Approved capability and implementation approach"),
        ("Implementation process", "Align\nValidate\nExpand"),
        ("Next steps", "Agree the next customer-owned implementation workshop"),
        (
            "Customer-safe legal notice",
            "Illustrative synthetic content only. Final scope is subject to an executed agreement.",
        ),
    ):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = heading
        slide.placeholders[1].text = body
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _textbox_only_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text = "Executive review"
    slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1)).text = "Jordan Lee, COO"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _pptx_with_unmapped_text() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Customer proposal"
    slide.placeholders[1].text = "Executive audience"
    slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.4)).text = "Stale customer name"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _limits() -> PptxLimits:
    return PptxLimits(
        max_bytes=50_000_000,
        max_slides=100,
        max_entries=2_000,
        max_expanded_bytes=250_000_000,
        max_media_assets=500,
        max_media_bytes=10_000_000,
        max_xml_bytes=5_000_000,
        max_extracted_characters=250_000,
    )


def _rewrite_pptx(
    source: bytes,
    *,
    additions: dict[str, bytes] | None = None,
    replacements: dict[str, bytes] | None = None,
) -> bytes:
    additions = additions or {}
    replacements = replacements or {}
    target = io.BytesIO()
    with (
        zipfile.ZipFile(io.BytesIO(source)) as original,
        zipfile.ZipFile(target, "w", compression=zipfile.ZIP_DEFLATED) as modified,
    ):
        for info in original.infolist():
            modified.writestr(info.filename, replacements.get(info.filename, original.read(info)))
        for name, content in additions.items():
            modified.writestr(name, content)
    return target.getvalue()


def _mark_first_zip_entry_encrypted(source: bytes) -> bytes:
    content = bytearray(source)
    for signature, flag_offset in ((b"PK\x03\x04", 6), (b"PK\x01\x02", 8)):
        position = content.find(signature)
        assert position >= 0
        flags = int.from_bytes(content[position + flag_offset : position + flag_offset + 2], "little") | 0x1
        content[position + flag_offset : position + flag_offset + 2] = flags.to_bytes(2, "little")
    return bytes(content)


def _run_worker(app: FastAPI) -> bool:
    engine = create_async_engine(TEST_DB_URL, connect_args={"check_same_thread": False})
    settings = app.state.settings

    async def run() -> bool:
        factory = async_sessionmaker(engine, expire_on_commit=False, autoflush=False)
        worker = CreateWorkerService(
            factory,
            settings,
            storage=create_visual_storage(settings),
            processor=create_processor(settings),
        )
        try:
            return await worker.run_once("create-test-worker")
        finally:
            await engine.dispose()

    return asyncio.run(run())


def _set_primary_membership_status(status: str) -> None:
    engine = create_async_engine(TEST_DB_URL, connect_args={"check_same_thread": False})

    async def change() -> None:
        factory = async_sessionmaker(engine, expire_on_commit=False)
        async with factory() as session:
            await session.execute(
                update(OrganisationMembership)
                .where(
                    OrganisationMembership.organisation_id == PRIMARY_ORGANISATION_ID,
                    OrganisationMembership.user_id == PRIMARY_USER_ID,
                )
                .values(status=status)
            )
            await session.commit()
        await engine.dispose()

    asyncio.run(change())


def _replace_approval_timestamp(version_id: str) -> None:
    engine = create_async_engine(TEST_DB_URL, connect_args={"check_same_thread": False})

    async def change() -> None:
        factory = async_sessionmaker(engine, expire_on_commit=False)
        async with factory() as session:
            await session.execute(
                update(CreatePresentationVersion)
                .where(CreatePresentationVersion.id == uuid.UUID(version_id))
                .values(approved_at=datetime.now(UTC) + timedelta(seconds=1))
            )
            await session.commit()
        await engine.dispose()

    asyncio.run(change())


def _expire_download_grant(token: str) -> None:
    engine = create_async_engine(TEST_DB_URL, connect_args={"check_same_thread": False})

    async def change() -> None:
        factory = async_sessionmaker(engine, expire_on_commit=False)
        async with factory() as session:
            expired_at = datetime.now(UTC) - timedelta(seconds=1)
            await session.execute(
                update(CreateDownloadGrant)
                .where(CreateDownloadGrant.token_hash == hashlib.sha256(token.encode()).hexdigest())
                .values(
                    created_at=expired_at - timedelta(seconds=1),
                    expires_at=expired_at,
                )
            )
            await session.commit()
        await engine.dispose()

    asyncio.run(change())


def _revoke_download_grant(token: str) -> None:
    engine = create_async_engine(TEST_DB_URL, connect_args={"check_same_thread": False})

    async def change() -> None:
        factory = async_sessionmaker(engine, expire_on_commit=False)
        async with factory() as session:
            await session.execute(
                update(CreateDownloadGrant)
                .where(CreateDownloadGrant.token_hash == hashlib.sha256(token.encode()).hexdigest())
                .values(revoked_at=datetime.now(UTC))
            )
            await session.commit()
        await engine.dispose()

    asyncio.run(change())


def _remove_expired_download_grants(app: FastAPI) -> dict[str, int]:
    engine = create_async_engine(TEST_DB_URL, connect_args={"check_same_thread": False})

    async def clean() -> dict[str, int]:
        factory = async_sessionmaker(engine, expire_on_commit=False)
        result = await run_retention(
            factory,
            app.state.settings,
            PRIMARY_ORGANISATION_ID,
            dry_run=False,
            batch_size=100,
        )
        await engine.dispose()
        return result.removed

    return asyncio.run(clean())


def _upload(client: TestClient) -> dict[str, object]:
    content = _pptx()
    response = client.post(
        "/api/v1/create/templates",
        json={
            "name": "Approved sales story",
            "fileName": "approved-sales-story.pptx",
            "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "contentBase64": base64.b64encode(content).decode(),
            "checksumSha256": hashlib.sha256(content).hexdigest(),
            "authorityAttested": True,
            "attestationVersion": 1,
        },
    )
    assert response.status_code == 202, response.text
    return response.json()


def _review_and_approve(client: TestClient, template: dict[str, object]) -> dict[str, object]:
    version = template["latestVersion"]
    assert isinstance(version, dict)
    for raw_slide in version["slides"]:
        assert isinstance(raw_slide, dict)
        is_locked_required = raw_slide["category"] == "appendix"
        is_title = raw_slide["category"] == "title"
        response = client.patch(
            f"/api/v1/create/template-slides/{raw_slide['id']}",
            json={
                "category": raw_slide["category"],
                "reuseState": "approved",
                "modificationPolicy": "locked" if is_locked_required else "text_placeholders_only",
                "customerSafe": True,
                "required": is_locked_required or is_title,
                "exactTextRequired": is_locked_required,
                "approvedDescription": "Approved for customer-facing reuse.",
                "placeholderMappings": {},
            },
        )
        assert response.status_code == 200, response.text
    response = client.post(
        f"/api/v1/create/templates/{template['id']}/versions/{version['id']}/approve",
        json={"confirmed": True},
    )
    assert response.status_code == 200, response.text
    return response.json()


def test_pptx_processor_rejects_active_content_and_external_relationships() -> None:
    processor = BoundedPptxProcessor(_limits())
    source = _pptx()
    parsed = processor.parse(source)
    assert len(parsed.slides) == 8

    unsafe = io.BytesIO()
    with (
        zipfile.ZipFile(io.BytesIO(source)) as original,
        zipfile.ZipFile(
            unsafe,
            "w",
            compression=zipfile.ZIP_DEFLATED,
        ) as modified,
    ):
        for info in original.infolist():
            modified.writestr(info.filename, original.read(info))
        modified.writestr("ppt/vbaProject.bin", b"not executable in RevenueOS")
    with pytest.raises(UnsafePptxError):
        processor.parse(unsafe.getvalue())

    rendered = processor.render(
        source,
        (RenderSlide(slide_number=2, replacements={}),),
        title="Customer-safe export",
        organisation_name="Must not be embedded",
    )
    exported = Presentation(io.BytesIO(rendered))
    assert len(exported.slides) == 1
    assert exported.core_properties.title == "Customer-safe export"
    assert exported.core_properties.author == "RevenueOS"
    assert exported.core_properties.last_modified_by == "RevenueOS"
    with zipfile.ZipFile(io.BytesIO(rendered)) as archive:
        content_types = archive.read("[Content_Types].xml")
        assert b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"' in content_types
        assert b"ns0:Types" not in content_types
        slide_parts = [
            name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        assert len(slide_parts) == 1
        xml_parts = {name: archive.read(name) for name in archive.namelist() if name.endswith(".xml")}
        package_text = b"\n".join(xml_parts.values())
    assert b"Product platform" not in package_text
    assert b"Implementation process" not in package_text
    assert b"Customer-safe legal notice" not in package_text
    assert b"TitlesOfParts" not in xml_parts["docProps/app.xml"]
    assert b"Must not be embedded" not in package_text


@pytest.mark.parametrize(
    "unsafe_path",
    [
        "../outside.xml",
        "ppt/activeX/activeX1.xml",
        "ppt/embeddings/oleObject1.bin",
        "ppt/fonts/font1.fntdata",
        "ppt/vbaProject.bin",
    ],
)
def test_pptx_processor_rejects_unsafe_package_paths(unsafe_path: str) -> None:
    with pytest.raises(UnsafePptxError):
        BoundedPptxProcessor(_limits()).parse(_rewrite_pptx(_pptx(), additions={unsafe_path: b"untrusted"}))


def test_pptx_processor_rejects_external_relationships_entities_and_unsafe_media() -> None:
    source = _pptx()
    rel_path = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(io.BytesIO(source)) as archive:
        relationships = archive.read(rel_path).replace(
            b"</Relationships>",
            (
                b'<Relationship Id="rExternal" '
                b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink" '
                b'Target="https://untrusted.example/" TargetMode="External"/></Relationships>'
            ),
        )
        presentation_xml = archive.read("ppt/presentation.xml")
    processor = BoundedPptxProcessor(_limits())
    with pytest.raises(UnsafePptxError):
        processor.parse(_rewrite_pptx(source, replacements={rel_path: relationships}))
    with pytest.raises(UnsafePptxError):
        processor.parse(
            _rewrite_pptx(
                source,
                replacements={
                    "ppt/presentation.xml": b'<!DOCTYPE p [<!ENTITY xxe SYSTEM "file:///etc/passwd">]>'
                    + presentation_xml
                },
            )
        )
    with pytest.raises(UnsafePptxError):
        processor.parse(
            _rewrite_pptx(
                source,
                replacements={
                    "ppt/presentation.xml": b" " * 5_000
                    + b'<!DOCTYPE p [<!ENTITY xxe SYSTEM "file:///etc/passwd">]>'
                    + presentation_xml
                },
            )
        )
    with pytest.raises(UnsupportedPptxError):
        processor.parse(_rewrite_pptx(source, additions={"ppt/media/image999.svg": b"<svg/>"}))


def test_pptx_processor_enforces_slide_media_and_decompression_limits() -> None:
    source = _pptx()
    strict_slides = replace(_limits(), max_slides=1)
    with pytest.raises(PptxLimitError):
        BoundedPptxProcessor(strict_slides).parse(source)
    strict_media = replace(_limits(), max_media_bytes=8)
    with pytest.raises(PptxLimitError):
        BoundedPptxProcessor(strict_media).parse(
            _rewrite_pptx(source, additions={"ppt/media/image999.png": b"\x89PNG\r\n\x1a\n0"})
        )
    with pytest.raises(PptxLimitError):
        BoundedPptxProcessor(_limits()).parse(_rewrite_pptx(source, additions={"ppt/unused.xml": b"x" * 1_100_000}))
    with pytest.raises(PptxLimitError):
        BoundedPptxProcessor(replace(_limits(), max_entries=10)).parse(source)
    with pytest.raises(PptxLimitError):
        BoundedPptxProcessor(replace(_limits(), max_expanded_bytes=1_000)).parse(source)
    with pytest.raises(UnsafePptxError):
        BoundedPptxProcessor(_limits()).parse(_mark_first_zip_entry_encrypted(source))
    malformed_directory = bytearray(source)
    central_directory = malformed_directory.find(b"PK\x01\x02")
    assert central_directory >= 0
    malformed_directory[central_directory : central_directory + 4] = b"BAD!"
    with pytest.raises(MalformedPptxError):
        BoundedPptxProcessor(_limits()).parse(bytes(malformed_directory))
    empty = io.BytesIO()
    Presentation().save(empty)
    with pytest.raises(MalformedPptxError):
        BoundedPptxProcessor(_limits()).parse(empty.getvalue())


def test_pptx_processor_marks_hidden_slides_unsafe_for_reuse_and_excludes_notes() -> None:
    presentation = Presentation(io.BytesIO(_pptx()))
    presentation.slides._sldIdLst[1].set("show", "0")  # noqa: SLF001 - fixture-only Open XML flag
    presentation.slides[0].notes_slide.notes_text_frame.text = "Internal presenter note"
    content = io.BytesIO()
    presentation.save(content)

    parsed = BoundedPptxProcessor(_limits()).parse(content.getvalue())
    assert parsed.slides[1].hidden is True
    assert parsed.slides[1].reuse_state == "excluded"
    assert parsed.slides[1].customer_safe is False
    assert "hidden_slides_excluded" in parsed.warning_codes
    assert "speaker_notes_excluded" in parsed.warning_codes
    rendered = BoundedPptxProcessor(_limits()).render(
        content.getvalue(),
        (RenderSlide(slide_number=1, replacements={}),),
        title="Customer-safe output",
        organisation_name="Not embedded",
    )
    with zipfile.ZipFile(io.BytesIO(rendered)) as archive:
        assert not any(name.startswith(("ppt/notesSlides/", "ppt/notesMasters/")) for name in archive.namelist())
        assert b"Internal presenter note" not in b"\n".join(
            archive.read(name) for name in archive.namelist() if name.endswith((".xml", ".rels"))
        )


@pytest.mark.parametrize(
    "unsafe_path",
    [
        "/absolute.xml",
        "C:/windows/path.xml",
        "ppt\\slides\\slide999.xml",
        "ppt/" + "a" * 252,
        "ppt/media/cafe\u0301.png",
    ],
)
def test_pptx_processor_rejects_ambiguous_or_non_canonical_paths(unsafe_path: str) -> None:
    with pytest.raises(UnsafePptxError):
        BoundedPptxProcessor(_limits()).parse(_rewrite_pptx(_pptx(), additions={unsafe_path: b"x"}))


def test_pptx_processor_rejects_duplicate_entries_polyglots_and_media_mismatch() -> None:
    source = _pptx()
    with zipfile.ZipFile(io.BytesIO(source)) as archive:
        duplicate_name = archive.namelist()[0]
        duplicate_data = archive.read(duplicate_name)
    with pytest.warns(UserWarning, match="Duplicate name"):
        duplicate = _rewrite_pptx(source, additions={duplicate_name: duplicate_data})
    processor = BoundedPptxProcessor(_limits())
    with pytest.raises(UnsafePptxError):
        processor.parse(duplicate)
    with pytest.raises(MalformedPptxError):
        processor.parse(source + b"untrusted trailing payload")
    with pytest.raises(UnsupportedPptxError):
        processor.parse(
            _rewrite_pptx(
                source,
                additions={"ppt/media/image999.jpg": b"\x89PNG\r\n\x1a\nnot-a-jpeg"},
            )
        )


def test_pptx_processor_rejects_missing_relationship_targets_and_deep_xml() -> None:
    source = _pptx()
    relationship_path = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(io.BytesIO(source)) as archive:
        relationships = archive.read(relationship_path).replace(
            b"../slideLayouts/slideLayout1.xml",
            b"../../../../outside.xml",
        )
    processor = BoundedPptxProcessor(_limits())
    with pytest.raises(UnsafePptxError):
        processor.parse(_rewrite_pptx(source, replacements={relationship_path: relationships}))
    deep_xml = (b"<root>" * 130) + (b"</root>" * 130)
    with pytest.raises(PptxLimitError):
        processor.parse(_rewrite_pptx(source, replacements={"ppt/presentation.xml": deep_xml}))
    relationship_xml = (
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        + "".join(
            f'<Relationship Id="rId{index}" Type="http://schemas.openxmlformats.org/officeDocument/2006/'
            'relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
            for index in range(1_001)
        )
        + "</Relationships>"
    ).encode()
    with pytest.raises(PptxLimitError):
        processor.parse(_rewrite_pptx(source, replacements={relationship_path: relationship_xml}))
    xinclude = b'<root xmlns:xi="http://www.w3.org/2001/XInclude"><xi:include href="file:///etc/passwd"/></root>'
    with pytest.raises(UnsafePptxError):
        processor.parse(_rewrite_pptx(source, replacements={"ppt/presentation.xml": xinclude}))


def test_pptx_renderer_writes_normal_unicode_as_editable_escaped_text() -> None:
    presentation = Presentation(io.BytesIO(_pptx()))
    title_shape = presentation.slides[0].shapes.title
    subtitle_shape = presentation.slides[0].placeholders[1]
    assert title_shape is not None
    title_shape.text_frame.paragraphs[0].font.name = "Aptos Display"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    source_buffer = io.BytesIO()
    presentation.save(source_buffer)
    source = source_buffer.getvalue()
    title = "Café & “Growth” <review> 上海 U0001f680"
    audience = "Zoë O’Connor"
    processor = BoundedPptxProcessor(_limits())
    rendered = processor.render(
        source,
        (
            RenderSlide(
                slide_number=1,
                replacements={title_shape.shape_id: title, subtitle_shape.shape_id: audience},
            ),
        ),
        title=title,
        organisation_name="Not embedded",
    )
    validated = processor.validate_output(
        rendered,
        (
            OutputSlideExpectation(
                title=title,
                body_blocks=(),
                exact_text_blocks=(),
                required=True,
                replacement_texts=(title, audience),
                replacement_shape_texts=(
                    (title_shape.shape_id, title),
                    (subtitle_shape.shape_id, audience),
                ),
            ),
        ),
    )
    assert validated.slide_count == 1
    reopened = Presentation(io.BytesIO(rendered))
    text = "\n".join(shape.text for shape in reopened.slides[0].shapes if shape.has_text_frame)
    assert title in text
    assert audience in text
    reopened_title = reopened.slides[0].shapes.title
    assert reopened_title is not None
    assert reopened_title.text_frame.paragraphs[0].font.name == "Aptos Display"
    assert reopened_title.text_frame.paragraphs[0].font.color.rgb == RGBColor(255, 255, 255)
    with pytest.raises(UnsafePptxError):
        processor.validate_output(
            rendered,
            (
                OutputSlideExpectation(
                    title=title,
                    body_blocks=(),
                    exact_text_blocks=(),
                    required=True,
                    replacement_shape_texts=((title_shape.shape_id, audience),),
                ),
            ),
        )


def test_generated_output_validator_fails_closed_on_review_mismatch_and_internal_ids() -> None:
    processor = BoundedPptxProcessor(_limits())
    rendered = processor.render(
        _pptx(),
        (RenderSlide(slide_number=2, replacements={}),),
        title="Validated output",
        organisation_name="Not embedded",
    )
    valid = processor.validate_output(
        rendered,
        (
            OutputSlideExpectation(
                title="Agenda",
                body_blocks=("Customer context",),
                exact_text_blocks=("Agenda",),
                required=True,
            ),
        ),
    )
    assert valid.slide_count == 1
    with pytest.raises(UnsafePptxError):
        processor.validate_output(
            rendered,
            (
                OutputSlideExpectation(
                    title="Agenda",
                    body_blocks=("A claim absent from the file",),
                    exact_text_blocks=(),
                    required=True,
                ),
            ),
        )
    with pytest.raises(UnsafePptxError):
        processor.validate_output(
            rendered,
            (
                OutputSlideExpectation(
                    title="Agenda",
                    body_blocks=(),
                    exact_text_blocks=(),
                    required=True,
                ),
            ),
            forbidden_values=("Agenda",),
        )


def test_claim_manifest_must_match_structured_slides_before_rendering() -> None:
    generated: list[dict[str, object]] = [
        {
            "planItemId": "plan-1",
            "title": "Customer priorities",
            "bodyBlocks": ["Approved customer statement"],
        }
    ]
    claims: list[dict[str, object]] = [
        {
            "planItemId": "plan-1",
            "claim": "Different statement",
            "reviewState": "not_required",
            "contentType": "customer_request",
        }
    ]
    with pytest.raises(GeneratedOutputValidationError):
        _validate_claim_manifest(generated, claims)


def test_template_upload_requires_authority_pptx_contract_and_valid_signature(client: TestClient) -> None:
    content = _pptx()
    base_payload = {
        "name": "Upload boundary",
        "fileName": "boundary.pptx",
        "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "contentBase64": base64.b64encode(content).decode(),
        "checksumSha256": hashlib.sha256(content).hexdigest(),
        "authorityAttested": True,
        "attestationVersion": 1,
    }
    wrong_extension = client.post(
        "/api/v1/create/templates",
        json={**base_payload, "fileName": "boundary.pdf"},
    )
    assert wrong_extension.status_code == 422
    missing_authority = client.post(
        "/api/v1/create/templates",
        json={**base_payload, "authorityAttested": False},
    )
    assert missing_authority.status_code == 422
    fake_content = b"PK this is not an Office package"
    renamed_non_pptx = client.post(
        "/api/v1/create/templates",
        json={
            **base_payload,
            "contentBase64": base64.b64encode(fake_content).decode(),
            "checksumSha256": hashlib.sha256(fake_content).hexdigest(),
        },
    )
    assert renamed_non_pptx.status_code == 422
    assert renamed_non_pptx.json()["code"] == "malformed_pptx"


def test_template_without_usable_placeholders_cannot_be_approved_as_editable(
    client: TestClient,
    app: FastAPI,
) -> None:
    content = _textbox_only_pptx()
    uploaded = client.post(
        "/api/v1/create/templates",
        json={
            "name": "Unmapped executive deck",
            "fileName": "unmapped-executive-deck.pptx",
            "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "contentBase64": base64.b64encode(content).decode(),
            "checksumSha256": hashlib.sha256(content).hexdigest(),
            "authorityAttested": True,
            "attestationVersion": 1,
        },
    )
    assert uploaded.status_code == 202, uploaded.text
    assert _run_worker(app)
    template = client.get(f"/api/v1/create/templates/{uploaded.json()['id']}").json()
    slide = template["latestVersion"]["slides"][0]
    editable = client.patch(
        f"/api/v1/create/template-slides/{slide['id']}",
        json={
            "category": "title",
            "reuseState": "approved",
            "modificationPolicy": "text_placeholders_only",
            "customerSafe": True,
            "required": True,
            "exactTextRequired": False,
            "approvedDescription": "Customer title slide.",
            "placeholderMappings": {},
        },
    )
    assert editable.status_code == 422
    assert editable.json()["code"] == "pptx_title_placeholders_required"

    reuse_only = client.patch(
        f"/api/v1/create/template-slides/{slide['id']}",
        json={
            "category": "title",
            "reuseState": "approved",
            "modificationPolicy": "reuse_as_is",
            "customerSafe": True,
            "required": True,
            "exactTextRequired": True,
            "approvedDescription": "Customer title slide.",
            "placeholderMappings": {},
        },
    )
    assert reuse_only.status_code == 200, reuse_only.text
    approval = client.post(
        f"/api/v1/create/templates/{template['id']}/versions/{template['latestVersion']['id']}/approve",
        json={"confirmed": True},
    )
    assert approval.status_code == 409
    assert approval.json()["code"] == "template_needs_attention"
    refreshed = client.get(f"/api/v1/create/templates/{template['id']}").json()
    assert refreshed["latestVersion"]["compatibilityState"] == "needs_attention"
    assert "pptx_title_placeholders_required" in refreshed["latestVersion"]["compatibilityDetails"]


def test_editable_slide_cannot_leave_unmapped_source_text_in_output(client: TestClient, app: FastAPI) -> None:
    content = _pptx_with_unmapped_text()
    uploaded = client.post(
        "/api/v1/create/templates",
        json={
            "name": "Stale source protection",
            "fileName": "stale-source-protection.pptx",
            "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "contentBase64": base64.b64encode(content).decode(),
            "checksumSha256": hashlib.sha256(content).hexdigest(),
            "authorityAttested": True,
            "attestationVersion": 1,
        },
    )
    assert uploaded.status_code == 202, uploaded.text
    assert _run_worker(app)
    template = client.get(f"/api/v1/create/templates/{uploaded.json()['id']}").json()
    slide = template["latestVersion"]["slides"][0]
    reviewed = client.patch(
        f"/api/v1/create/template-slides/{slide['id']}",
        json={
            "category": "title",
            "reuseState": "approved",
            "modificationPolicy": "text_placeholders_only",
            "customerSafe": True,
            "required": True,
            "exactTextRequired": False,
            "approvedDescription": "Customer title slide.",
            "placeholderMappings": {},
        },
    )
    assert reviewed.status_code == 422
    assert reviewed.json()["code"] == "pptx_unmapped_text_requires_lock"


def test_create_template_plan_generation_review_approval_and_private_download(
    client: TestClient,
    app: FastAPI,
    caplog: pytest.LogCaptureFixture,
) -> None:
    availability = client.get("/api/v1/create/availability")
    assert availability.status_code == 200
    assert availability.json()["state"] == "available"

    queued = _upload(client)
    assert queued["latestVersion"]["processingState"] == "processing"  # type: ignore[index]
    assert _run_worker(app)
    template = client.get(f"/api/v1/create/templates/{queued['id']}").json()
    assert template["latestVersion"]["processingState"] == "ready"
    assert template["latestVersion"]["slideCount"] == 8
    template = _review_and_approve(client, template)
    assert template["latestVersion"]["approvalState"] == "approved"
    assert template["latestVersion"]["compatibilityState"] == "compatible"

    company = create_company(client, name="Create Customer")
    opportunity = create_opportunity(client, str(company["id"]), name="Secure rollout")
    contact = create_contact(client, str(company["id"]), first_name="Taylor")
    created = client.post(
        "/api/v1/create/presentations",
        json={
            "accountId": company["id"],
            "opportunityId": opportunity["id"],
            "objective": "solution_overview",
            "audience": [
                {
                    "contactId": contact["id"],
                    "audienceType": "executive",
                }
            ],
            "templateVersionId": template["latestVersion"]["id"],
            "focusInstruction": "Keep the implementation discussion concise.",
            "idempotencyKey": "create-plan-1",
        },
    )
    assert created.status_code == 201, created.text
    presentation = created.json()
    assert presentation["state"] == "draft_plan"
    assert presentation["accountName"] == "Create Customer"
    assert presentation["plan"][0]["required"] is True
    assert any(
        item["category"] == "appendix" and item["required"] and item["exactTextRequired"]
        for item in presentation["plan"]
    )

    generated = client.post(
        f"/api/v1/create/presentations/{presentation['id']}/generate",
        json={"idempotencyKey": "create-generate-1"},
    )
    assert generated.status_code == 200, generated.text
    assert generated.json()["state"] == "generating"
    assert _run_worker(app)
    presentation = client.get(f"/api/v1/create/presentations/{presentation['id']}").json()
    assert presentation["state"] == "needs_review"
    assert presentation["currentVersion"]["slides"]
    assert all(claim["origin"] == "approved_company_content" for claim in presentation["currentVersion"]["claims"])

    editable = next(
        slide
        for slide in presentation["currentVersion"]["slides"]
        if slide["modificationPolicy"] == "text_placeholders_only" and slide["title"] != presentation["title"]
    )
    blocked = client.patch(
        f"/api/v1/create/presentations/{presentation['id']}/slides/{editable['planItemId']}",
        json={
            "title": "Internal recommendation",
            "bodyBlocks": ["Manager coaching says the win probability is 75%."],
        },
    )
    assert blocked.status_code == 422
    assert blocked.json()["code"] == "internal_only_content"

    edited = client.patch(
        f"/api/v1/create/presentations/{presentation['id']}/slides/{editable['planItemId']}",
        json={
            "title": editable["title"],
            "bodyBlocks": ["Confirm the customer-owned implementation workshop."],
        },
    )
    assert edited.status_code == 200, edited.text
    assert edited.json()["state"] == "generating"
    assert _run_worker(app)
    presentation = client.get(f"/api/v1/create/presentations/{presentation['id']}").json()
    edited_claim = next(
        claim for claim in presentation["currentVersion"]["claims"] if claim["planItemId"] == editable["planItemId"]
    )
    assert edited_claim["origin"] == "user_edited"
    assert edited_claim["reviewState"] == "pending"
    reviewed = client.post(
        f"/api/v1/create/presentations/{presentation['id']}/review",
        json={"decisions": [{"claimId": edited_claim["id"], "action": "keep"}]},
    )
    assert reviewed.status_code == 200, reviewed.text
    approved = client.post(
        f"/api/v1/create/presentations/{presentation['id']}/approve",
        json={"confirmed": True},
    )
    assert approved.status_code == 200, approved.text
    assert approved.json()["currentVersion"]["downloadAvailable"] is True

    grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant")
    assert grant.status_code == 200, grant.text
    assert "?" not in grant.json()["downloadUrl"]
    assert grant.headers["cache-control"] == "private, no-store"
    assert grant.headers["referrer-policy"] == "no-referrer"
    downloaded = client.post(
        grant.json()["downloadUrl"],
        json={"grantToken": grant.json()["grantToken"]},
    )
    assert downloaded.status_code == 200, downloaded.text
    assert downloaded.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert downloaded.headers["cache-control"] == "private, no-store"
    assert downloaded.headers["referrer-policy"] == "no-referrer"
    assert downloaded.headers["x-content-type-options"] == "nosniff"
    assert downloaded.headers["content-disposition"].startswith('attachment; filename="')
    assert downloaded.headers["content-disposition"].endswith('.pptx"')
    assert "\r" not in downloaded.headers["content-disposition"]
    assert "\n" not in downloaded.headers["content-disposition"]
    rendered = Presentation(io.BytesIO(downloaded.content))
    rendered_text = "\n".join(shape.text for slide in rendered.slides for shape in slide.shapes if shape.has_text_frame)
    assert "Confirm the customer-owned implementation workshop." in rendered_text
    assert "Taylor" in rendered_text
    assert "win probability" not in rendered_text.casefold()
    assert rendered.core_properties.title == presentation["title"]
    assert rendered.core_properties.comments == ""
    replay = client.post(
        grant.json()["downloadUrl"],
        json={"grantToken": grant.json()["grantToken"]},
    )
    assert replay.status_code == 403
    assert replay.json()["code"] == "invalid_download_grant"
    assert grant.json()["grantToken"] not in caplog.text

    expired_grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant").json()
    _expire_download_grant(expired_grant["grantToken"])
    expired_download = client.post(
        expired_grant["downloadUrl"],
        json={"grantToken": expired_grant["grantToken"]},
    )
    assert expired_download.status_code == 403
    assert expired_download.json()["code"] == "invalid_download_grant"
    assert _remove_expired_download_grants(app)["expired_create_download_grants"] == 1

    revoked_grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant").json()
    _revoke_download_grant(revoked_grant["grantToken"])
    revoked_download = client.post(
        revoked_grant["downloadUrl"],
        json={"grantToken": revoked_grant["grantToken"]},
    )
    assert revoked_download.status_code == 403
    assert revoked_download.json()["code"] == "invalid_download_grant"

    cross_tenant_grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant").json()
    app.dependency_overrides[get_current_user] = cast_auth_dependency(secondary_user())
    try:
        cross_tenant_download = client.post(
            cross_tenant_grant["downloadUrl"],
            json={"grantToken": cross_tenant_grant["grantToken"]},
        )
        assert cross_tenant_download.status_code == 404
        assert cross_tenant_download.json()["code"] == "presentation_not_found"
    finally:
        app.dependency_overrides.pop(get_current_user)

    revoked_membership_grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant").json()
    _set_primary_membership_status("disabled")
    try:
        revoked_membership_download = client.post(
            revoked_membership_grant["downloadUrl"],
            json={"grantToken": revoked_membership_grant["grantToken"]},
        )
        assert revoked_membership_download.status_code == 403
    finally:
        _set_primary_membership_status("active")

    approval_bound_grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant").json()
    _replace_approval_timestamp(presentation["currentVersion"]["id"])
    approval_changed = client.post(
        approval_bound_grant["downloadUrl"],
        json={"grantToken": approval_bound_grant["grantToken"]},
    )
    assert approval_changed.status_code == 403
    assert approval_changed.json()["code"] == "invalid_download_grant"

    integrity_grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant").json()
    stored_file = next(TEST_VISUAL_STORAGE.rglob(f"{presentation['currentVersion']['id']}.pptx"))
    expected_content = stored_file.read_bytes()
    missing_file = stored_file.with_suffix(".missing")
    stored_file.rename(missing_file)
    try:
        unavailable = client.post(
            integrity_grant["downloadUrl"],
            json={"grantToken": integrity_grant["grantToken"]},
        )
        assert unavailable.status_code == 409
        assert unavailable.json()["code"] == "presentation_file_unavailable"
    finally:
        missing_file.rename(stored_file)

    integrity_grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant").json()
    stored_file.write_bytes(b"corrupt")
    try:
        integrity_failure = client.post(
            integrity_grant["downloadUrl"],
            json={"grantToken": integrity_grant["grantToken"]},
        )
        assert integrity_failure.status_code == 409
        assert integrity_failure.json()["code"] == "presentation_file_integrity_failed"
    finally:
        stored_file.write_bytes(expected_content)

    concurrent_grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant").json()

    def use_concurrent_grant() -> int:
        with TestClient(app) as concurrent_client:
            return concurrent_client.post(
                concurrent_grant["downloadUrl"],
                json={"grantToken": concurrent_grant["grantToken"]},
            ).status_code

    with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
        statuses = sorted(executor.map(lambda _: use_concurrent_grant(), range(2)))
    assert statuses == [200, 403]

    app.state.settings.private_beta_max_create_presentations_per_user_per_day = 1
    quota_denied = client.post(
        f"/api/v1/create/presentations/{presentation['id']}/generate",
        json={
            "idempotencyKey": "create-generate-over-user-limit",
            "explicitRegenerate": True,
        },
    )
    assert quota_denied.status_code == 429
    assert quota_denied.json()["code"] == "create_user_daily_limit"
    app.state.settings.private_beta_max_create_presentations_per_user_per_day = 10
    app.state.settings.private_beta_max_create_presentations_per_organisation_per_day = 1
    organisation_quota_denied = client.post(
        f"/api/v1/create/presentations/{presentation['id']}/generate",
        json={
            "idempotencyKey": "create-generate-over-organisation-limit",
            "explicitRegenerate": True,
        },
    )
    assert organisation_quota_denied.status_code == 429
    assert organisation_quota_denied.json()["code"] == "create_organisation_daily_limit"

    member = AuthenticatedUser(
        user_id=PRIMARY_USER_ID,
        external_auth_id="user_dev_001",
        display_name="Create Member",
        email="member@example.test",
        organisation_id=PRIMARY_ORGANISATION_ID,
        organisation_name="Example Revenue Team",
        organisation_slug="example-revenue-team",
        role="member",
        auth_mode="mock",
    )
    app.dependency_overrides[get_current_user] = cast_auth_dependency(member)
    assert client.get(f"/api/v1/create/templates/{template['id']}").status_code == 200
    source = _pptx()
    member_upload = client.post(
        "/api/v1/create/templates",
        json={
            "name": "Member cannot upload",
            "fileName": "member.pptx",
            "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "contentBase64": base64.b64encode(source).decode(),
            "checksumSha256": hashlib.sha256(source).hexdigest(),
            "authorityAttested": True,
            "attestationVersion": 1,
        },
    )
    assert member_upload.status_code == 403
    member_plan = client.post(
        "/api/v1/create/presentations",
        json={
            "accountId": company["id"],
            "opportunityId": opportunity["id"],
            "objective": "executive_presentation",
            "audience": [{"name": "Operations leadership", "audienceType": "executive"}],
            "templateVersionId": template["latestVersion"]["id"],
            "idempotencyKey": "create-member-plan-1",
        },
    )
    assert member_plan.status_code == 201, member_plan.text


def test_approved_business_case_flows_into_create_with_exact_scenario_provenance(
    client: TestClient,
    app: FastAPI,
) -> None:
    company = create_company(client, name="Business Case Customer")
    opportunity = create_opportunity(client, str(company["id"]), name="Access transformation")
    model = _create_approved_model(client)
    created_case = client.post(
        "/api/v1/create/business-cases",
        json={
            "accountId": company["id"],
            "opportunityId": opportunity["id"],
            "modelVersionId": model["latestVersion"]["id"],
            "currency": "AUD",
            "idempotencyKey": "create-integration-case-1",
        },
    )
    assert created_case.status_code == 201, created_case.text
    business_case = created_case.json()
    calculated = client.post(
        f"/api/v1/create/business-cases/{business_case['id']}/calculate",
        json={
            "inputs": _inputs(),
            "scenarios": [
                {"name": "conservative", "overrides": [{"key": "minutes_future", "value": "10"}]},
                {"name": "upside", "overrides": [{"key": "minutes_future", "value": "2"}]},
            ],
            "idempotencyKey": "create-integration-calculate-1",
        },
    )
    assert calculated.status_code == 200, calculated.text
    approved_case = client.post(
        f"/api/v1/create/business-cases/{business_case['id']}/approve",
        json={"confirmed": True},
    )
    assert approved_case.status_code == 200, approved_case.text
    case_version = approved_case.json()["currentVersion"]

    queued = _upload(client)
    assert _run_worker(app)
    template = client.get(f"/api/v1/create/templates/{queued['id']}").json()
    template = _review_and_approve(client, template)
    presentation_response = client.post(
        "/api/v1/create/presentations",
        json={
            "accountId": company["id"],
            "opportunityId": opportunity["id"],
            "objective": "business_case",
            "audience": [{"name": "Finance leadership", "audienceType": "finance"}],
            "templateVersionId": template["latestVersion"]["id"],
            "businessCaseVersionId": case_version["id"],
            "businessCaseScenario": "all",
            "idempotencyKey": "create-business-case-plan-1",
        },
    )
    assert presentation_response.status_code == 201, presentation_response.text
    presentation = presentation_response.json()
    assert presentation["businessCaseId"] == business_case["id"]
    assert presentation["businessCaseVersionId"] == case_version["id"]
    assert presentation["businessCaseScenario"] == "all"
    assert any("approved_business_case" in item["sourceClasses"] for item in presentation["plan"])

    generated = client.post(
        f"/api/v1/create/presentations/{presentation['id']}/generate",
        json={"idempotencyKey": "create-business-case-generate-1"},
    )
    assert generated.status_code == 200, generated.text
    assert _run_worker(app)
    presentation = client.get(f"/api/v1/create/presentations/{presentation['id']}").json()
    business_claims = [
        claim for claim in presentation["currentVersion"]["claims"] if claim["origin"] == "approved_business_case"
    ]
    assert business_claims
    assert all(claim["sourceIds"] == [case_version["id"]] for claim in business_claims)
    claim_text = " ".join(claim["claim"] for claim in business_claims)
    assert "conservative assumptions" in claim_text
    assert "base-case assumptions" in claim_text
    assert "upside assumptions" in claim_text
    assert "Material assumption" in claim_text
    assert "not a guarantee of future results" in claim_text
    assert "safe_divide" not in claim_text
    approved_presentation = client.post(
        f"/api/v1/create/presentations/{presentation['id']}/approve",
        json={"confirmed": True},
    )
    assert approved_presentation.status_code == 200, approved_presentation.text
    case_grant = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant")
    assert case_grant.status_code == 200, case_grant.text
    case_download = client.post(
        case_grant.json()["downloadUrl"],
        json={"grantToken": case_grant.json()["grantToken"]},
    )
    assert case_download.status_code == 200, case_download.text
    case_deck = Presentation(io.BytesIO(case_download.content))
    case_deck_text = "\n".join(
        shape.text for slide in case_deck.slides for shape in slide.shapes if shape.has_text_frame
    )
    assert all(claim["claim"] in case_deck_text for claim in business_claims)

    superseded = client.post(
        f"/api/v1/create/business-cases/{business_case['id']}/calculate",
        json={"inputs": _inputs(rekey_cost="0"), "idempotencyKey": "create-integration-calculate-2"},
    )
    assert superseded.status_code == 200, superseded.text
    blocked_export = client.post(f"/api/v1/create/presentations/{presentation['id']}/download-grant")
    assert blocked_export.status_code == 409
    assert blocked_export.json()["code"] == "claim_source_changed"

    app.dependency_overrides[get_current_user] = cast_auth_dependency(secondary_user())
    cross_tenant = client.get(f"/api/v1/create/templates/{template['id']}")
    assert cross_tenant.status_code == 404
    assert cross_tenant.json()["code"] == "template_not_found"


def test_create_entitlement_fails_closed_and_is_admin_managed(client: TestClient) -> None:
    disabled = client.patch("/api/v1/create/admin/entitlement", json={"enabled": False})
    assert disabled.status_code == 200
    assert disabled.json()["state"] == "not_in_plan"
    denied = client.get("/api/v1/create/templates")
    assert denied.status_code == 403
    assert denied.json()["code"] == "create_not_entitled"
    enabled = client.patch("/api/v1/create/admin/entitlement", json={"enabled": True})
    assert enabled.status_code == 200
    assert enabled.json()["state"] == "available"
