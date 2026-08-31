"""Generate the synthetic WO-039B PPTX compatibility evidence set."""

from __future__ import annotations

import base64
import hashlib
import io
import json
import sys
import zipfile
from dataclasses import asdict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

REPOSITORY_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPOSITORY_ROOT / "apps" / "api" / "src"))

from revenueos.create_pptx import (  # noqa: E402
    CREATE_PPTX_PROFILE_VERSION,
    BoundedPptxProcessor,
    OutputSlideExpectation,
    PptxLimits,
    PptxProcessingError,
    RenderSlide,
)

EVIDENCE_ROOT = REPOSITORY_ROOT / "docs" / "07-sprints" / "assets" / "wo-039b"
SOURCE_ROOT = EVIDENCE_ROOT / "source"
GENERATED_ROOT = EVIDENCE_ROOT / "generated"
SYNTHETIC_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9YEksEIAAAAASUVORK5CYII="
)


def limits() -> PptxLimits:
    return PptxLimits(
        max_bytes=50_000_000,
        max_slides=100,
        max_entries=2_000,
        max_expanded_bytes=250_000_000,
        max_media_assets=500,
        max_media_bytes=10_000_000,
        max_xml_bytes=5_000_000,
        max_extracted_characters=250_000,
    )


def presentation_bytes(presentation: Presentation) -> bytes:
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def add_title_slide(presentation: Presentation, title: str, subtitle: str) -> tuple[int, int]:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    assert slide.shapes.title is not None
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide.shapes.title.shape_id, slide.placeholders[1].shape_id


def add_content_slide(presentation: Presentation, title: str, body: str) -> tuple[int, int]:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    assert slide.shapes.title is not None
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    return slide.shapes.title.shape_id, slide.placeholders[1].shape_id


def simple_corporate() -> tuple[Presentation, tuple[RenderSlide, ...], tuple[OutputSlideExpectation, ...]]:
    deck = Presentation()
    title_id, audience_id = add_title_slide(deck, "Northstar corporate story", "Approved source audience")
    body_title_id, body_id = add_content_slide(deck, "Customer priorities", "Approved source statement")
    replacements = (
        RenderSlide(1, {title_id: "Harbour Health — Executive review", audience_id: "Jordan Lee, COO"}),
        RenderSlide(
            2,
            {
                body_title_id: "Customer priorities",
                body_id: "Reduce manual access reviews\nKeep implementation customer-owned",
            },
        ),
    )
    expectations = (
        OutputSlideExpectation(
            "Harbour Health — Executive review",
            (),
            (),
            True,
            replacement_shape_texts=tuple(replacements[0].replacements.items()),
        ),
        OutputSlideExpectation(
            "Customer priorities",
            ("Reduce manual access reviews", "Keep implementation customer-owned"),
            (),
            True,
            replacement_shape_texts=tuple(replacements[1].replacements.items()),
        ),
    )
    return deck, replacements, expectations


def brand_heavy() -> tuple[Presentation, tuple[RenderSlide, ...], tuple[OutputSlideExpectation, ...]]:
    deck = Presentation()
    title_id, audience_id = add_title_slide(deck, "Summit brand", "Approved source audience")
    title_slide = deck.slides[0]
    background = title_slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(13, 30, 55)
    brand_bar = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), deck.slide_height)
    brand_bar.fill.solid()
    brand_bar.fill.fore_color.rgb = RGBColor(45, 212, 191)
    brand_bar.line.fill.background()
    title_slide.shapes.add_picture(io.BytesIO(SYNTHETIC_PNG), Inches(11.7), Inches(0.4), Inches(0.6))
    assert title_slide.shapes.title is not None
    for paragraph in title_slide.shapes.title.text_frame.paragraphs:
        paragraph.font.name = "Aptos Display"
        paragraph.font.size = Pt(30)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
    body_title_id, body_id = add_content_slide(deck, "Approved capability", "Approved source content")
    deck.slides[1].shapes.add_picture(io.BytesIO(SYNTHETIC_PNG), Inches(11.6), Inches(6.4), Inches(0.45))
    replacements = (
        RenderSlide(1, {title_id: "Atlas Operations — Solution overview", audience_id: "Executive team"}),
        RenderSlide(
            2,
            {
                body_title_id: "Approved capability",
                body_id: "Centralised access administration\nPolicy-based customer controls",
            },
        ),
    )
    expectations = (
        OutputSlideExpectation(
            "Atlas Operations — Solution overview",
            (),
            (),
            True,
            replacement_shape_texts=tuple(replacements[0].replacements.items()),
        ),
        OutputSlideExpectation(
            "Approved capability",
            ("Centralised access administration", "Policy-based customer controls"),
            (),
            True,
            replacement_shape_texts=tuple(replacements[1].replacements.items()),
        ),
    )
    return deck, replacements, expectations


def multi_layout() -> tuple[Presentation, tuple[RenderSlide, ...], tuple[OutputSlideExpectation, ...]]:
    deck = Presentation()
    title_id, audience_id = add_title_slide(deck, "Multi-layout source", "Approved audience")
    section = deck.slides.add_slide(deck.slide_layouts[2])
    assert section.shapes.title is not None
    section.shapes.title.text = "Delivery approach"
    section.placeholders[1].text = "Approved section introduction"
    section_title_id = section.shapes.title.shape_id
    section_body_id = section.placeholders[1].shape_id
    content_title_id, content_body_id = add_content_slide(deck, "Next steps", "Approved next step")
    replacements = (
        RenderSlide(1, {title_id: "BluePeak — Technical workshop", audience_id: "Platform leadership"}),
        RenderSlide(
            2,
            {
                section_title_id: "Delivery approach",
                section_body_id: "Align\nValidate\nExpand",
            },
        ),
        RenderSlide(
            3,
            {
                content_title_id: "Next steps",
                content_body_id: "Confirm the customer-owned implementation workshop",
            },
        ),
    )
    expectations = (
        OutputSlideExpectation(
            "BluePeak — Technical workshop",
            (),
            (),
            True,
            replacement_shape_texts=tuple(replacements[0].replacements.items()),
        ),
        OutputSlideExpectation(
            "Delivery approach",
            ("Align", "Validate", "Expand"),
            (),
            True,
            replacement_shape_texts=tuple(replacements[1].replacements.items()),
        ),
        OutputSlideExpectation(
            "Next steps",
            ("Confirm the customer-owned implementation workshop",),
            (),
            True,
            replacement_shape_texts=tuple(replacements[2].replacements.items()),
        ),
    )
    return deck, replacements, expectations


def exact_legal() -> tuple[Presentation, tuple[RenderSlide, ...], tuple[OutputSlideExpectation, ...]]:
    deck = Presentation()
    title_id, audience_id = add_title_slide(deck, "Legal source", "Approved audience")
    legal_title_id, legal_body_id = add_content_slide(
        deck,
        "Customer-safe legal notice",
        "Illustrative synthetic content only. Final scope is subject to an executed agreement.",
    )
    deck.slides[1].notes_slide.notes_text_frame.text = "INTERNAL: legal review routing must never ship"
    legal_title = "Customer-safe legal notice"
    legal_body = "Illustrative synthetic content only. Final scope is subject to an executed agreement."
    replacements = (
        RenderSlide(1, {title_id: "Cedar Labs — Proposal", audience_id: "Commercial leadership"}),
        RenderSlide(2, {}),
    )
    expectations = (
        OutputSlideExpectation(
            "Cedar Labs — Proposal",
            (),
            (),
            True,
            replacement_shape_texts=tuple(replacements[0].replacements.items()),
        ),
        OutputSlideExpectation(legal_title, (legal_body,), (legal_title, legal_body), True),
    )
    assert legal_title_id > 0 and legal_body_id > 0
    return deck, replacements, expectations


def add_external_relationship(content: bytes) -> bytes:
    relationship_path = "ppt/slides/_rels/slide1.xml.rels"
    target = io.BytesIO()
    with (
        zipfile.ZipFile(io.BytesIO(content)) as source,
        zipfile.ZipFile(target, "w", compression=zipfile.ZIP_DEFLATED) as output,
    ):
        for info in source.infolist():
            data = source.read(info)
            if info.filename == relationship_path:
                data = data.replace(
                    b"</Relationships>",
                    b'<Relationship Id="rUnsafe" Type="http://schemas.openxmlformats.org/officeDocument/2006/'
                    b'relationships/hyperlink" Target="https://invalid.example/" TargetMode="External"/>'
                    b"</Relationships>",
                )
            output.writestr(info.filename, data)
    return target.getvalue()


def package_summary(content: bytes) -> dict[str, object]:
    with zipfile.ZipFile(io.BytesIO(content)) as archive:
        names = archive.namelist()
        xml = b"\n".join(archive.read(name) for name in names if name.endswith((".xml", ".rels")))
    deck = Presentation(io.BytesIO(content))
    return {
        "sha256": hashlib.sha256(content).hexdigest(),
        "byteSize": len(content),
        "slideCount": len(deck.slides),
        "slideText": [
            [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
            for slide in deck.slides
        ],
        "editableTextShapeCount": sum(1 for slide in deck.slides for shape in slide.shapes if shape.has_text_frame),
        "hasNotes": any(name.startswith(("ppt/notesSlides/", "ppt/notesMasters/")) for name in names),
        "hasComments": any(name.startswith("ppt/comments/") for name in names),
        "hasExternalRelationship": b'TargetMode="External"' in xml,
        "hasInternalEvidenceMarker": b"INTERNAL:" in xml,
    }


def main() -> None:
    processor = BoundedPptxProcessor(limits())
    results: list[dict[str, object]] = []
    for name, factory in (
        ("simple-corporate", simple_corporate),
        ("brand-heavy", brand_heavy),
        ("multi-layout", multi_layout),
        ("exact-legal", exact_legal),
    ):
        deck, replacements, expectations = factory()
        source = presentation_bytes(deck)
        source_path = SOURCE_ROOT / f"{name}-source.pptx"
        generated_path = GENERATED_ROOT / f"{name}-generated.pptx"
        source_path.write_bytes(source)
        parsed = processor.parse(source)
        generated = processor.render(
            source,
            replacements,
            title=expectations[0].title,
            organisation_name="Synthetic evidence organisation",
        )
        validation = processor.validate_output(
            generated,
            expectations,
            forbidden_values=("INTERNAL: legal review routing must never ship",),
        )
        generated_path.write_bytes(generated)
        results.append(
            {
                "fixture": name,
                "classification": "PASS",
                "profileVersion": CREATE_PPTX_PROFILE_VERSION,
                "source": package_summary(source),
                "sourceWarningCodes": list(parsed.warning_codes),
                "generated": package_summary(generated),
                "outputValidation": asdict(validation),
                "structuralGuarantee": "Generated PPTX reparsed and matched expected replacement/exact text.",
                "visualLimitation": "Fonts, wrapping and layout remain renderer- and installed-font-dependent.",
            }
        )

    unsupported_source = add_external_relationship(presentation_bytes(simple_corporate()[0]))
    unsupported_path = SOURCE_ROOT / "unsupported-external-content-source.pptx"
    unsupported_path.write_bytes(unsupported_source)
    failure_code = "unexpected_success"
    try:
        processor.parse(unsupported_source)
    except PptxProcessingError as exc:
        failure_code = exc.code
    results.append(
        {
            "fixture": "unsupported-external-content",
            "classification": "UNSUPPORTED",
            "profileVersion": CREATE_PPTX_PROFILE_VERSION,
            "source": {
                "sha256": hashlib.sha256(unsupported_source).hexdigest(),
                "byteSize": len(unsupported_source),
            },
            "safeFailureCode": failure_code,
            "structuralGuarantee": "External relationships are rejected before python-pptx processing.",
            "visualLimitation": "Rejected fixtures are deliberately not rendered.",
        }
    )
    if failure_code != "unsafe_pptx":
        raise RuntimeError(f"Unsupported fixture did not fail safely: {failure_code}")
    (EVIDENCE_ROOT / "inspection-manifest.json").write_text(
        json.dumps(
            {
                "schemaVersion": 1,
                "syntheticDataOnly": True,
                "profileVersion": CREATE_PPTX_PROFILE_VERSION,
                "fixtures": results,
            },
            indent=2,
            sort_keys=True,
        )
        + "\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
